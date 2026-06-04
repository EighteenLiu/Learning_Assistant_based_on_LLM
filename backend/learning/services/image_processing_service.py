from __future__ import annotations

from pathlib import Path
import shutil
from typing import Dict, List

from django.conf import settings
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.util import Pt
from PIL import Image, ImageDraw, ImageFont

from app.debug_logger import debug_log

from .ppt_parser_service import PPTParserService
from .slide_render_service import SlideRenderService


class ImageProcessingService:
    EMU_PER_INCH = 914400
    PIXELS_PER_INCH = 96
    MIN_MARGIN_PT = 0
    CONTRAST_THRESHOLD = 3.2
    LONG_PAGE_CHAR_THRESHOLD = max(
        int(getattr(settings, "TRANSLATION_LONG_PAGE_CHAR_THRESHOLD", 1800) or 200),
        200,
    )
    LONG_PAGE_MAX_FONT_SIZE = max(
        int(getattr(settings, "TRANSLATION_LONG_PAGE_MAX_FONT_SIZE", 9) or 6),
        6,
    )

    @staticmethod
    # 实现课件内容提取，把文本、位置和样式转成后续可用的结构。
    def _parse_shape_path(shape_path: str) -> list[int]:
        if not shape_path:
            return []
        return [int(item) for item in str(shape_path).split(".") if item != ""]

    @staticmethod
    # 实现 _find_shape_by_path 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _find_shape_by_path(shapes, shape_path: str):
        indices = ImageProcessingService._parse_shape_path(shape_path)
        current_shapes = shapes
        current_shape = None
        for depth, index in enumerate(indices):
            if index < 0 or index >= len(current_shapes):
                return None
            current_shape = current_shapes[index]
            if depth < len(indices) - 1:
                current_shapes = getattr(current_shape, "shapes", None)
                if current_shapes is None:
                    return None
        return current_shape

    @staticmethod
    # 实现 _contains_hangul 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _contains_hangul(text: str) -> bool:
        return any("\uac00" <= char <= "\ud7af" for char in str(text or ""))

    @staticmethod
    # 实现 _preferred_font_info 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _preferred_font_info(text_frame, paragraphs: list[str] | None = None) -> tuple[str, int]:
        font_name = "Microsoft YaHei"
        max_size = 20

        for paragraph in getattr(text_frame, "paragraphs", []):
            for run in getattr(paragraph, "runs", []):
                font = getattr(run, "font", None)
                if font is None:
                    continue
                if getattr(font, "name", None):
                    font_name = font.name
                if getattr(font, "size", None):
                    try:
                        max_size = max(max_size, int(font.size.pt))
                    except Exception:
                        pass
                if font_name and max_size:
                    break
            if font_name and max_size:
                break

        text_samples = "\n".join(paragraphs or [])
        if ImageProcessingService._contains_hangul(text_samples):
            font_name = "Malgun Gothic"
        return font_name, max_size

    @staticmethod
    # 实现 _original_font_info 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _original_font_info(
        text_frame,
        paragraphs: list[str] | None = None,
        font_name_hint: str | None = None,
        font_size_hint: float | int | None = None,
    ) -> tuple[str, int]:
        fallback_name, fallback_size = ImageProcessingService._preferred_font_info(text_frame, paragraphs)
        resolved_name = str(font_name_hint or fallback_name or "Microsoft YaHei").strip()

        resolved_size = fallback_size
        if font_size_hint:
            try:
                resolved_size = max(4, int(round(float(font_size_hint))))
            except Exception:
                resolved_size = fallback_size

        return resolved_name or "Microsoft YaHei", resolved_size

    @staticmethod
    # 实现 _preferred_font_rgb 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _preferred_font_rgb(text_frame) -> tuple[int, int, int] | None:
        candidates = []
        for paragraph in getattr(text_frame, "paragraphs", []):
            if getattr(paragraph.font, "color", None) is not None:
                candidates.append(paragraph.font.color)
            for run in getattr(paragraph, "runs", []):
                font = getattr(run, "font", None)
                if font is not None and getattr(font, "color", None) is not None:
                    candidates.append(font.color)

        for color in candidates:
            try:
                rgb = getattr(color, "rgb", None)
                if rgb is not None:
                    return int(rgb[0]), int(rgb[1]), int(rgb[2])
            except Exception:
                continue
        return None

    @staticmethod
    # 实现 _font_candidates 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _font_candidates(font_name: str, text: str = "") -> list[str]:
        lowered = str(font_name or "").lower()
        ordered = []
        if "malgun" in lowered or "맑은" in lowered or ImageProcessingService._contains_hangul(text):
            ordered.extend(
                [
                    "C:\\Windows\\Fonts\\malgun.ttf",
                    "C:\\Windows\\Fonts\\malgunsl.ttf",
                    "C:\\Windows\\Fonts\\gulim.ttc",
                ]
            )
        if "yahei" in lowered or "微软雅黑" in lowered:
            ordered.append("C:\\Windows\\Fonts\\msyh.ttc")
        if "heiti" in lowered or "黑体" in lowered:
            ordered.append("C:\\Windows\\Fonts\\simhei.ttf")
        if "song" in lowered or "宋体" in lowered:
            ordered.append("C:\\Windows\\Fonts\\simsun.ttc")
        ordered.extend(
            [
                "C:\\Windows\\Fonts\\msyh.ttc",
                "C:\\Windows\\Fonts\\simhei.ttf",
                "C:\\Windows\\Fonts\\simsun.ttc",
                "C:\\Windows\\Fonts\\malgun.ttf",
                "C:\\Windows\\Fonts\\malgunsl.ttf",
                "C:\\Windows\\Fonts\\gulim.ttc",
                "C:\\Windows\\Fonts\\arial.ttf",
            ]
        )
        deduped: list[str] = []
        for path in ordered:
            if path not in deduped and Path(path).exists():
                deduped.append(path)
        return deduped

    @staticmethod
    # 实现 _load_font 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _load_font(font_name: str, font_size: int, text: str = "") -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
        for font_path in ImageProcessingService._font_candidates(font_name, text):
            try:
                return ImageFont.truetype(font_path, max(font_size, 1))
            except Exception:
                continue
        return ImageFont.load_default()

    @staticmethod
    # 实现 _emu_to_px 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _emu_to_px(value: int) -> int:
        return max(int((value / ImageProcessingService.EMU_PER_INCH) * ImageProcessingService.PIXELS_PER_INCH), 1)

    @staticmethod
    # 实现 _wrap_line 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _wrap_line(font: ImageFont.FreeTypeFont | ImageFont.ImageFont, text: str, max_width: int) -> list[str]:
        if not text:
            return [""]

        lines: list[str] = []
        current = ""
        for char in text:
            candidate = f"{current}{char}"
            bbox = font.getbbox(candidate or " ")
            width = bbox[2] - bbox[0]
            if current and width > max_width:
                lines.append(current.rstrip())
                current = char.lstrip()
            else:
                current = candidate
        if current:
            lines.append(current.rstrip())
        return lines or [text]

    @staticmethod
    # 实现 _measure_paragraphs 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _measure_paragraphs(
        font: ImageFont.FreeTypeFont | ImageFont.ImageFont,
        paragraphs: list[str],
        max_width: int,
        line_spacing_ratio: float = 0.32,
        paragraph_spacing_px: int = 0,
    ) -> tuple[int, int]:
        wrapped_groups: list[list[str]] = []
        for paragraph in paragraphs:
            wrapped_groups.append(ImageProcessingService._wrap_line(font, paragraph, max_width))

        spacing = max(1, int(getattr(font, "size", 12) * line_spacing_ratio))
        widths: list[int] = []
        total_height = 0
        for paragraph_index, wrapped_lines in enumerate(wrapped_groups):
            for index, line in enumerate(wrapped_lines):
                bbox = font.getbbox(line or " ")
                widths.append(bbox[2] - bbox[0])
                total_height += bbox[3] - bbox[1]
                if index < len(wrapped_lines) - 1:
                    total_height += spacing
            if paragraph_index < len(wrapped_groups) - 1:
                total_height += paragraph_spacing_px
        return max(widths or [0]), max(total_height, 1)

    @staticmethod
    # 实现 _available_box_px 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _available_box_px(
        text_frame,
        width_emu: int,
        height_emu: int,
        margin_pt: float | None = None,
    ) -> tuple[int, int]:
        margin_left = int(getattr(text_frame, "margin_left", 0) or 0)
        margin_right = int(getattr(text_frame, "margin_right", 0) or 0)
        margin_top = int(getattr(text_frame, "margin_top", 0) or 0)
        margin_bottom = int(getattr(text_frame, "margin_bottom", 0) or 0)
        if margin_pt is not None:
            margin_emu = int(Pt(margin_pt))
            margin_left = margin_right = margin_top = margin_bottom = margin_emu

        width_emu -= margin_left + margin_right
        height_emu -= margin_top + margin_bottom
        return (
            ImageProcessingService._emu_to_px(max(width_emu, 1)),
            ImageProcessingService._emu_to_px(max(height_emu, 1)),
        )

    @staticmethod
    # 实现 _best_fit_font_size 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _best_fit_font_size(
        text_frame,
        paragraphs: list[str],
        font_name: str,
        max_size: int,
        min_size: int = 6,
        width_emu: int | None = None,
        height_emu: int | None = None,
        line_spacing_ratio: float = 0.32,
        paragraph_spacing_px: int = 0,
        margin_pt: float | None = None,
    ) -> tuple[int, bool]:
        # 通过二分查找字体大小，比逐号递减更快；导出大课件时这个函数会被频繁调用。
        width_emu = int(width_emu or 0)
        height_emu = int(height_emu or 0)
        if not width_emu or not height_emu:
            try:
                parent = getattr(text_frame, "_parent", None)
                width_emu = width_emu or int(getattr(parent, "width", 0) or 0)
                height_emu = height_emu or int(getattr(parent, "height", 0) or 0)
            except Exception:
                pass

        available_width, available_height = ImageProcessingService._available_box_px(
            text_frame,
            width_emu,
            height_emu,
            margin_pt=margin_pt,
        )
        low = min_size
        high = max(max_size, min_size)
        best = min_size
        fits = False
        text_sample = "\n".join(paragraphs)

        while low <= high:
            mid = (low + high) // 2
            font = ImageProcessingService._load_font(font_name, mid, text_sample)
            text_width, text_height = ImageProcessingService._measure_paragraphs(
                font,
                paragraphs,
                available_width,
                line_spacing_ratio=line_spacing_ratio,
                paragraph_spacing_px=paragraph_spacing_px,
            )
            if text_width <= available_width and text_height <= available_height:
                best = mid
                fits = True
                low = mid + 1
            else:
                high = mid - 1

        if not fits:
            font = ImageProcessingService._load_font(font_name, best, text_sample)
            text_width, text_height = ImageProcessingService._measure_paragraphs(
                font,
                paragraphs,
                available_width,
                line_spacing_ratio=line_spacing_ratio,
                paragraph_spacing_px=paragraph_spacing_px,
            )
            fits = text_width <= available_width and text_height <= available_height
        return best, fits

    @staticmethod
    # 实现 _apply_font_size 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _apply_font_size(
        text_frame,
        font_name: str,
        font_size: int,
        line_spacing_ratio: float = 0.32,
        margin_pt: float | None = None,
        paragraph_spacing_pt: float = 0,
    ) -> None:
        if margin_pt is not None:
            margin_emu = int(Pt(margin_pt))
            try:
                text_frame.margin_left = margin_emu
                text_frame.margin_right = margin_emu
                text_frame.margin_top = margin_emu
                text_frame.margin_bottom = margin_emu
            except Exception:
                pass

        for paragraph in getattr(text_frame, "paragraphs", []):
            try:
                paragraph.font.name = font_name
                paragraph.font.size = Pt(font_size)
                paragraph.space_before = Pt(0)
                paragraph.space_after = Pt(max(paragraph_spacing_pt, 0))
                paragraph.line_spacing = max(0.72, 1 + line_spacing_ratio)
            except Exception:
                pass
            runs = getattr(paragraph, "runs", [])
            for run in runs:
                try:
                    run.font.name = font_name
                    run.font.size = Pt(font_size)
                except Exception:
                    pass

    @staticmethod
    # 实现 _relative_luminance 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _relative_luminance(rgb: tuple[int, int, int]) -> float:
        channels = []
        for channel in rgb:
            value = max(0, min(255, int(channel))) / 255
            if value <= 0.03928:
                channels.append(value / 12.92)
            else:
                channels.append(((value + 0.055) / 1.055) ** 2.4)
        return 0.2126 * channels[0] + 0.7152 * channels[1] + 0.0722 * channels[2]

    @staticmethod
    # 实现 _contrast_ratio 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _contrast_ratio(foreground_rgb: tuple[int, int, int], background_rgb: tuple[int, int, int]) -> float:
        fg = ImageProcessingService._relative_luminance(foreground_rgb)
        bg = ImageProcessingService._relative_luminance(background_rgb)
        lighter = max(fg, bg)
        darker = min(fg, bg)
        return (lighter + 0.05) / (darker + 0.05)

    @staticmethod
    # 实现 _shape_background_rgb 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _shape_background_rgb(host) -> tuple[int, int, int] | None:
        fill = getattr(host, "fill", None)
        if fill is None:
            return None
        try:
            if getattr(fill, "transparency", 0) >= 1:
                return None
        except Exception:
            pass
        try:
            rgb = getattr(fill.fore_color, "rgb", None)
            if rgb is not None:
                return int(rgb[0]), int(rgb[1]), int(rgb[2])
        except Exception:
            return None
        return None

    @staticmethod
    # 实现 _best_contrast_rgb 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _best_contrast_rgb(
        source_rgb: tuple[int, int, int] | None,
        background_rgb: tuple[int, int, int] | None,
    ) -> tuple[int, int, int] | None:
        if background_rgb is None:
            return source_rgb

        # 译文会覆盖到原始背景上，必要时自动换成高对比色，保证导出的课件可读。
        if source_rgb and ImageProcessingService._contrast_ratio(source_rgb, background_rgb) >= ImageProcessingService.CONTRAST_THRESHOLD:
            return source_rgb

        candidates = [
            (15, 23, 42),
            (248, 250, 252),
            (30, 41, 59),
            (255, 255, 255),
            (0, 0, 0),
        ]
        ranked = sorted(
            candidates,
            key=lambda rgb: ImageProcessingService._contrast_ratio(rgb, background_rgb),
            reverse=True,
        )
        return ranked[0] if ranked else source_rgb

    @staticmethod
    # 实现 _apply_font_color 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _apply_font_color(text_frame, rgb: tuple[int, int, int] | None) -> None:
        if rgb is None:
            return
        color = RGBColor(*rgb)
        for paragraph in getattr(text_frame, "paragraphs", []):
            try:
                paragraph.font.color.rgb = color
            except Exception:
                pass
            for run in getattr(paragraph, "runs", []):
                try:
                    run.font.color.rgb = color
                except Exception:
                    pass

    @staticmethod
    # 实现 _resolve_media_path 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _resolve_media_path(media_url: str) -> Path | None:
        normalized = str(media_url or "").strip()
        if not normalized:
            return None
        media_prefix = str(getattr(settings, "MEDIA_URL", "") or "")
        if media_prefix and normalized.startswith(media_prefix):
            normalized = normalized[len(media_prefix) :]
        normalized = normalized.lstrip("/\\")
        path = Path(settings.MEDIA_ROOT) / Path(normalized)
        return path if path.exists() else None

    @staticmethod
    # 实现 _sample_background_rgb 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _sample_background_rgb(source_image: Image.Image | None, container: dict) -> tuple[int, int, int] | None:
        if source_image is None:
            return None

        image_width, image_height = source_image.size
        left = max(int(float(container.get("x", 0) or 0) * image_width), 0)
        top = max(int(float(container.get("y", 0) or 0) * image_height), 0)
        width = max(int(float(container.get("w", 0) or 0) * image_width), 1)
        height = max(int(float(container.get("h", 0) or 0) * image_height), 1)
        right = min(left + width, image_width)
        bottom = min(top + height, image_height)
        if right <= left or bottom <= top:
            return None

        try:
            sample = source_image.crop((left, top, right, bottom)).convert("RGB").resize((1, 1))
            return tuple(int(channel) for channel in sample.getpixel((0, 0)))
        except Exception:
            return None

    @staticmethod
    # 实现 _page_text_char_count 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _page_text_char_count(containers: list[dict]) -> int:
        total = 0
        for container in containers:
            text = str(container.get("translated_text") or container.get("text") or "").strip()
            if not text:
                continue
            total += len(text)
        return total

    @staticmethod
    # 实现 _set_text_frame_content 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _set_text_frame_content(
        text_frame,
        paragraphs: list[str],
        width_emu: int | None = None,
        height_emu: int | None = None,
        host=None,
        background_rgb: tuple[int, int, int] | None = None,
        font_name_hint: str | None = None,
        font_size_hint: float | int | None = None,
        force_compact: bool = False,
    ) -> None:
        # 这里集中处理字体、换行、缩放和颜色，避免 PPT 文本框、表格单元格、覆盖层各自实现一套规则。
        normalized_paragraphs = [str(item) for item in paragraphs]
        combined_text = "\n".join(normalized_paragraphs).strip()
        font_name, max_size = ImageProcessingService._original_font_info(
            text_frame,
            normalized_paragraphs,
            font_name_hint=font_name_hint,
            font_size_hint=font_size_hint,
        )
        preferred_rgb = ImageProcessingService._preferred_font_rgb(text_frame)
        dense_hangul = ImageProcessingService._contains_hangul(combined_text) and len(combined_text) >= 40
        text_frame.text = combined_text
        text_frame.word_wrap = True
        try:
            text_frame.auto_size = MSO_AUTO_SIZE.NONE
        except Exception:
            pass

        if force_compact:
            # 文字特别密集的页面优先保证完整显示，因此主动降低字号和边距。
            max_size = min(max_size, ImageProcessingService.LONG_PAGE_MAX_FONT_SIZE)
            layout_candidates = [
                {"line_spacing_ratio": 0.02, "margin_pt": ImageProcessingService.MIN_MARGIN_PT, "paragraph_spacing_pt": 0, "min_size": 4, "auto_fit": True},
            ]
        elif dense_hangul:
            layout_candidates = [
                {"line_spacing_ratio": 0.1, "margin_pt": 1, "paragraph_spacing_pt": 0, "min_size": 5, "auto_fit": False},
                {"line_spacing_ratio": 0.02, "margin_pt": ImageProcessingService.MIN_MARGIN_PT, "paragraph_spacing_pt": 0, "min_size": 4, "auto_fit": True},
            ]
        else:
            layout_candidates = [
                {"line_spacing_ratio": 0.18, "margin_pt": None, "paragraph_spacing_pt": 0, "min_size": 6, "auto_fit": False},
                {"line_spacing_ratio": 0.08, "margin_pt": 1, "paragraph_spacing_pt": 0, "min_size": 5, "auto_fit": False},
                {"line_spacing_ratio": 0.02, "margin_pt": ImageProcessingService.MIN_MARGIN_PT, "paragraph_spacing_pt": 0, "min_size": 4, "auto_fit": True},
            ]
        selected_font_size = 6
        selected_candidate = layout_candidates[-1]
        fits_any_layout = False

        for candidate in layout_candidates:
            best_size, fits = ImageProcessingService._best_fit_font_size(
                text_frame,
                normalized_paragraphs,
                font_name,
                max_size,
                min_size=candidate["min_size"],
                width_emu=width_emu,
                height_emu=height_emu,
                line_spacing_ratio=candidate["line_spacing_ratio"],
                margin_pt=candidate["margin_pt"],
            )
            selected_font_size = best_size
            selected_candidate = candidate
            if fits:
                fits_any_layout = True
                break

        ImageProcessingService._apply_font_size(
            text_frame,
            font_name,
            selected_font_size,
            line_spacing_ratio=selected_candidate["line_spacing_ratio"],
            margin_pt=selected_candidate["margin_pt"],
            paragraph_spacing_pt=selected_candidate["paragraph_spacing_pt"],
        )

        effective_background = background_rgb or ImageProcessingService._shape_background_rgb(host) or (255, 255, 255)
        adjusted_rgb = ImageProcessingService._best_contrast_rgb(preferred_rgb, effective_background)
        ImageProcessingService._apply_font_color(text_frame, adjusted_rgb)

        try:
            text_frame.fit_text(font_family=font_name, max_size=selected_font_size)
        except Exception:
            pass

        if not fits_any_layout and selected_candidate.get("auto_fit"):
            try:
                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass
        if force_compact:
            try:
                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass

    @staticmethod
    # 实现 _apply_container_translation 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _apply_container_translation(
        slide,
        container: dict,
        background_rgb: tuple[int, int, int] | None = None,
        force_compact: bool = False,
    ) -> bool:
        translated_paragraphs = container.get("translated_paragraphs") or container.get("paragraphs") or []
        translated_paragraphs = [str(item).strip() for item in translated_paragraphs]
        translated_text = str(container.get("translated_text") or container.get("text") or "").strip()
        if not any(translated_paragraphs) and translated_text:
            translated_paragraphs = [translated_text]
        if not translated_paragraphs:
            return False

        try:
            kind = str(container.get("kind", "") or "")
            shape = ImageProcessingService._find_shape_by_path(slide.shapes, str(container.get("shape_path", "")))
            if kind == "image_ocr":
                if shape is None:
                    return False
                overlay = slide.shapes.add_textbox(
                    int(getattr(shape, "left", 0) or 0),
                    int(getattr(shape, "top", 0) or 0),
                    int(getattr(shape, "width", 0) or 0),
                    int(getattr(shape, "height", 0) or 0),
                )
                ImageProcessingService._set_text_frame_content(
                    overlay.text_frame,
                    translated_paragraphs,
                    width_emu=int(getattr(overlay, "width", 0) or 0),
                    height_emu=int(getattr(overlay, "height", 0) or 0),
                    host=overlay,
                    background_rgb=background_rgb,
                    font_name_hint=str(container.get("font_name", "") or ""),
                    font_size_hint=container.get("font_size_pt"),
                    force_compact=force_compact,
                )
                return True

            if shape is None:
                return False

            if kind == "table_cell":
                row_index = int(container.get("row_index") or 0)
                col_index = int(container.get("col_index") or 0)
                cell = shape.table.cell(row_index, col_index)
                cell_width = int(getattr(shape.table.columns[col_index], "width", 0) or 0)
                cell_height = int(getattr(shape.table.rows[row_index], "height", 0) or 0)
                ImageProcessingService._set_text_frame_content(
                    cell.text_frame,
                    translated_paragraphs,
                    width_emu=cell_width,
                    height_emu=cell_height,
                    host=cell,
                    background_rgb=background_rgb,
                    font_name_hint=str(container.get("font_name", "") or ""),
                    font_size_hint=container.get("font_size_pt"),
                    force_compact=force_compact,
                )
                return True

            if getattr(shape, "has_text_frame", False):
                ImageProcessingService._set_text_frame_content(
                    shape.text_frame,
                    translated_paragraphs,
                    width_emu=int(getattr(shape, "width", 0) or 0),
                    height_emu=int(getattr(shape, "height", 0) or 0),
                    host=shape,
                    background_rgb=background_rgb,
                    font_name_hint=str(container.get("font_name", "") or ""),
                    font_size_hint=container.get("font_size_pt"),
                    force_compact=force_compact,
                )
                return True
        except Exception as exc:
            debug_log(
                hypothesisId="H11",
                runId="pre-diagnose",
                location="image_processing_service:_apply_container_translation",
                message="Failed to apply translated text to shape",
                data={
                    "shape_path": str(container.get("shape_path", "")),
                    "shape_id": container.get("shape_id"),
                    "error": str(exc)[:300],
                },
            )
        return False

    @staticmethod
    # 实现翻译处理步骤，负责组织输入、调用模型并整理译文结果。
    def translated_pptx_path(courseware_id: int) -> Path:
        export_dir = Path(settings.MEDIA_ROOT) / "translated_pptx"
        export_dir.mkdir(parents=True, exist_ok=True)
        return export_dir / f"courseware_{courseware_id}_translated.pptx"

    @staticmethod
    # 实现翻译处理步骤，负责组织输入、调用模型并整理译文结果。
    def translated_pdf_path(courseware_id: int) -> Path:
        export_dir = Path(settings.MEDIA_ROOT) / "translated_pdf"
        export_dir.mkdir(parents=True, exist_ok=True)
        return export_dir / f"courseware_{courseware_id}_translated.pdf"

    @staticmethod
    # 实现翻译处理步骤，负责组织输入、调用模型并整理译文结果。
    def translated_output_path(courseware_id: int, source_file_path: str) -> Path:
        if PPTParserService.is_pdf_file(source_file_path):
            return ImageProcessingService.translated_pdf_path(courseware_id)
        return ImageProcessingService.translated_pptx_path(courseware_id)

    @staticmethod
    # 实现 _prepare_output_dir 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _prepare_output_dir(courseware_id: int, output_folder: str) -> Path:
        output_dir = Path(settings.MEDIA_ROOT) / output_folder / f"courseware_{courseware_id}"
        if output_dir.exists():
            shutil.rmtree(output_dir, ignore_errors=True)
        output_dir.mkdir(parents=True, exist_ok=True)
        return output_dir

    @staticmethod
    # 实现 _container_pixel_rect 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _container_pixel_rect(source_image: Image.Image, container: dict) -> tuple[int, int, int, int]:
        image_width, image_height = source_image.size
        left = max(int(float(container.get("x", 0) or 0) * image_width), 0)
        top = max(int(float(container.get("y", 0) or 0) * image_height), 0)
        width = max(int(float(container.get("w", 0) or 0) * image_width), 1)
        height = max(int(float(container.get("h", 0) or 0) * image_height), 1)
        right = min(left + width, image_width)
        bottom = min(top + height, image_height)
        return left, top, right, bottom

    @staticmethod
    # 实现课件预览或导出处理，把布局数据转换为可视化结果。
    def _best_fit_font_size_for_image(
        paragraphs: list[str],
        font_name: str,
        max_size: int,
        min_size: int,
        max_width: int,
        max_height: int,
        line_spacing_ratio: float = 0.16,
        paragraph_spacing_px: int = 0,
    ) -> tuple[int, bool]:
        low = max(min_size, 1)
        high = max(max_size, low)
        best = low
        fits = False
        text_sample = "\n".join(paragraphs)

        while low <= high:
            mid = (low + high) // 2
            font = ImageProcessingService._load_font(font_name, mid, text_sample)
            text_width, text_height = ImageProcessingService._measure_paragraphs(
                font,
                paragraphs,
                max(max_width, 1),
                line_spacing_ratio=line_spacing_ratio,
                paragraph_spacing_px=paragraph_spacing_px,
            )
            if text_width <= max_width and text_height <= max_height:
                best = mid
                fits = True
                low = mid + 1
            else:
                high = mid - 1
        return best, fits

    @staticmethod
    # 实现课件预览或导出处理，把布局数据转换为可视化结果。
    def _draw_paragraphs_on_image(
        source_image: Image.Image,
        container: dict,
        *,
        force_compact_page: bool = False,
    ) -> bool:
        translated_paragraphs = container.get("translated_paragraphs") or container.get("paragraphs") or []
        translated_paragraphs = [str(item).strip() for item in translated_paragraphs if str(item).strip()]
        translated_text = str(container.get("translated_text") or container.get("text") or "").strip()
        if not translated_paragraphs and translated_text:
            translated_paragraphs = [translated_text]
        if not translated_paragraphs:
            return False

        left, top, right, bottom = ImageProcessingService._container_pixel_rect(source_image, container)
        if right <= left or bottom <= top:
            return False

        box_width = max(right - left, 1)
        box_height = max(bottom - top, 1)
        font_name = str(container.get("font_name", "") or "Microsoft YaHei")
        size_hint = 0
        try:
            size_hint = int(round(float(container.get("font_size_pt") or 0)))
        except Exception:
            size_hint = 0
        max_size = size_hint if size_hint > 0 else min(max(int(box_height * 0.7), 9), 48)
        min_size = 4 if force_compact_page else 6
        if force_compact_page:
            max_size = min(max_size, ImageProcessingService.LONG_PAGE_MAX_FONT_SIZE)
        max_size = max(max_size, min_size)

        inner_padding_x = max(int(box_width * 0.03), 1)
        inner_padding_y = max(int(box_height * 0.03), 1)
        available_width = max(box_width - inner_padding_x * 2, 1)
        available_height = max(box_height - inner_padding_y * 2, 1)
        line_spacing_ratio = 0.08 if force_compact_page else 0.16

        font_size, _ = ImageProcessingService._best_fit_font_size_for_image(
            translated_paragraphs,
            font_name,
            max_size=max_size,
            min_size=min_size,
            max_width=available_width,
            max_height=available_height,
            line_spacing_ratio=line_spacing_ratio,
            paragraph_spacing_px=0,
        )
        font = ImageProcessingService._load_font(font_name, font_size, "\n".join(translated_paragraphs))

        background_rgb = ImageProcessingService._sample_background_rgb(source_image, container) or (255, 255, 255)
        text_rgb = ImageProcessingService._best_contrast_rgb(None, background_rgb) or (15, 23, 42)
        draw = ImageDraw.Draw(source_image)
        draw.rectangle((left, top, right, bottom), fill=background_rgb)

        cursor_y = top + inner_padding_y
        line_spacing = max(int(font_size * line_spacing_ratio), 1)
        for paragraph_index, paragraph in enumerate(translated_paragraphs):
            wrapped_lines = ImageProcessingService._wrap_line(font, paragraph, available_width)
            for line in wrapped_lines:
                if cursor_y >= bottom:
                    return True
                draw.text((left + inner_padding_x, cursor_y), line, fill=text_rgb, font=font)
                bbox = font.getbbox(line or " ")
                line_height = max(int(bbox[3] - bbox[1]), 1)
                cursor_y += line_height + line_spacing
            if paragraph_index < len(translated_paragraphs) - 1:
                cursor_y += max(1, int(font_size * 0.1))
        return True

    @staticmethod
    # 实现翻译处理步骤，负责组织输入、调用模型并整理译文结果。
    def _render_translated_images_from_source(
        courseware_id: int,
        slides_data: List[Dict],
        output_folder: str,
    ) -> tuple[Dict[int, str], dict[int, Path]]:
        output_dir = ImageProcessingService._prepare_output_dir(courseware_id, output_folder)
        media_url_prefix = f"{settings.MEDIA_URL}{output_folder}/courseware_{courseware_id}"
        generated_urls: Dict[int, str] = {}
        generated_paths: dict[int, Path] = {}

        for slide_data in sorted(slides_data, key=lambda item: int(item.get("slide_no") or 0)):
            # PDF 或不可编辑页面用“原图 + 译文覆盖”的方式生成预览，尽量保留原页面视觉效果。
            slide_no = int(slide_data.get("slide_no") or 0)
            if slide_no <= 0:
                continue

            source_image_path = ImageProcessingService._resolve_media_path(str(slide_data.get("source_image_url") or ""))
            if source_image_path is None or not source_image_path.exists():
                continue

            translated_layout = slide_data.get("translated_layout", {}) or {}
            translated_containers = translated_layout.get("text_containers", []) or []
            page_text_chars = ImageProcessingService._page_text_char_count(translated_containers)
            force_compact_page = page_text_chars >= ImageProcessingService.LONG_PAGE_CHAR_THRESHOLD

            target_name = source_image_path.name or f"幻灯片{slide_no}.PNG"
            target_path = output_dir / target_name
            try:
                with Image.open(source_image_path) as image:
                    rendered = image.convert("RGB")
                    for container in translated_containers:
                        ImageProcessingService._draw_paragraphs_on_image(
                            rendered,
                            container,
                            force_compact_page=force_compact_page,
                        )
                    rendered.save(target_path, format="PNG")
            except Exception as exc:
                debug_log(
                    hypothesisId="H13",
                    runId="pre-diagnose",
                    location="image_processing_service:_render_translated_images_from_source",
                    message="Failed to render translated page image",
                    data={"courseware_id": courseware_id, "slide_no": slide_no, "error": str(exc)[:300]},
                )
                continue

            generated_urls[slide_no] = f"{media_url_prefix}/{target_path.name}"
            generated_paths[slide_no] = target_path

        return generated_urls, generated_paths

    @staticmethod
    # 实现 _apply_layouts_to_presentation 对应的核心处理，封装输入转换、状态更新或结果返回。
    def _apply_layouts_to_presentation(presentation: Presentation, slides_data: List[Dict]) -> None:
        for slide_data in slides_data:
            slide_no = int(slide_data.get("slide_no") or 0)
            if slide_no <= 0 or slide_no > len(presentation.slides):
                continue

            translated_layout = slide_data.get("translated_layout", {}) or {}
            translated_containers = translated_layout.get("text_containers", []) or []
            slide = presentation.slides[slide_no - 1]
            page_text_chars = ImageProcessingService._page_text_char_count(translated_containers)
            force_compact_page = page_text_chars >= ImageProcessingService.LONG_PAGE_CHAR_THRESHOLD
            source_image = None
            source_image_path = ImageProcessingService._resolve_media_path(str(slide_data.get("source_image_url") or ""))
            if source_image_path:
                try:
                    source_image = Image.open(source_image_path).convert("RGB")
                except Exception:
                    source_image = None

            for container in translated_containers:
                # 优先回写原 shape/table；如果是图片 OCR 区域，则用覆盖层放置译文。
                background_rgb = ImageProcessingService._sample_background_rgb(source_image, container)
                ImageProcessingService._apply_container_translation(
                    slide,
                    container,
                    background_rgb=background_rgb,
                    force_compact=force_compact_page,
                )

            if source_image is not None:
                try:
                    source_image.close()
                except Exception:
                    pass

    @staticmethod
    # 实现翻译处理步骤，负责组织输入、调用模型并整理译文结果。
    def export_translated_courseware(courseware_id: int, slides_data: List[Dict], source_file_path: str = "") -> Path | None:
        if PPTParserService.is_pdf_file(source_file_path):
            return ImageProcessingService.export_translated_pdf(courseware_id, slides_data, source_file_path)
        return ImageProcessingService.export_translated_pptx(courseware_id, slides_data, source_file_path)

    @staticmethod
    # 实现翻译处理步骤，负责组织输入、调用模型并整理译文结果。
    def export_translated_pptx(courseware_id: int, slides_data: List[Dict], source_ppt_path: str = "") -> Path | None:
        if not source_ppt_path:
            return None

        temp_source_path: Path | None = None
        output_path = ImageProcessingService.translated_pptx_path(courseware_id)
        try:
            editable_path, temp_source_path = PPTParserService.ensure_editable_pptx(source_ppt_path)
            presentation = Presentation(editable_path)
            ImageProcessingService._apply_layouts_to_presentation(presentation, slides_data)
            presentation.save(str(output_path))
            return output_path
        except Exception as exc:
            debug_log(
                hypothesisId="H13",
                runId="pre-diagnose",
                location="image_processing_service:export_translated_pptx",
                message="Failed to export translated PPTX",
                data={"courseware_id": courseware_id, "error": str(exc)[:500]},
            )
            return None
        finally:
            if temp_source_path:
                try:
                    temp_source_path.unlink(missing_ok=True)
                except Exception:
                    pass

    @staticmethod
    # 实现翻译处理步骤，负责组织输入、调用模型并整理译文结果。
    def export_translated_pdf(courseware_id: int, slides_data: List[Dict], source_pdf_path: str = "") -> Path | None:
        if not source_pdf_path:
            return None
        try:
            import fitz  # type: ignore
        except Exception as exc:
            raise RuntimeError("PDF export requires PyMuPDF. Please install it with `pip install PyMuPDF`.") from exc

        if not PPTParserService.is_pdf_file(source_pdf_path):
            return None

        # PDF 导出无法像 PPTX 一样编辑对象，因此先渲染每页图片，再把译文覆盖后的图片重新组装成 PDF。
        SlideRenderService.export_slide_images(source_pdf_path, courseware_id, output_folder="rendered_slides")
        _, translated_image_paths = ImageProcessingService._render_translated_images_from_source(
            courseware_id,
            slides_data,
            output_folder="translated_pdf_pages",
        )
        output_path = ImageProcessingService.translated_pdf_path(courseware_id)

        source_doc = None
        target_doc = None
        try:
            source_doc = fitz.open(str(Path(source_pdf_path).resolve()))
            target_doc = fitz.open()
            render_dpi = max(int(getattr(settings, "PDF_RENDER_DPI", 144) or 72), 72)
            matrix = fitz.Matrix(max(float(render_dpi) / 72.0, 1.0), max(float(render_dpi) / 72.0, 1.0))
            for page_index, page in enumerate(source_doc, start=1):
                rect = page.rect
                target_page = target_doc.new_page(width=rect.width, height=rect.height)
                translated_image_path = translated_image_paths.get(page_index)
                if translated_image_path is not None and translated_image_path.exists():
                    target_page.insert_image(rect, filename=str(translated_image_path))
                    continue

                pixmap = page.get_pixmap(matrix=matrix, alpha=False)
                target_page.insert_image(rect, stream=pixmap.tobytes("png"))
            target_doc.save(str(output_path), deflate=True)
            return output_path
        except Exception as exc:
            debug_log(
                hypothesisId="H13",
                runId="pre-diagnose",
                location="image_processing_service:export_translated_pdf",
                message="Failed to export translated PDF",
                data={"courseware_id": courseware_id, "error": str(exc)[:500]},
            )
            return None
        finally:
            if source_doc is not None:
                try:
                    source_doc.close()
                except Exception:
                    pass
            if target_doc is not None:
                try:
                    target_doc.close()
                except Exception:
                    pass

    @staticmethod
    # 实现课件预览或导出处理，把布局数据转换为可视化结果。
    def process_all_slides(courseware_id: int, slides_data: List[Dict], source_ppt_path: str = "") -> Dict[int, str]:
        source_path = str(source_ppt_path or "")
        if PPTParserService.is_pdf_file(source_path):
            # PDF 预览直接走图片覆盖链路，不依赖 PowerPoint 或 LibreOffice。
            SlideRenderService.export_slide_images(source_path, courseware_id, output_folder="rendered_slides")
            processed_urls, _ = ImageProcessingService._render_translated_images_from_source(
                courseware_id,
                slides_data,
                output_folder="processed_slides",
            )
            return processed_urls

        output_pptx = ImageProcessingService.export_translated_pptx(courseware_id, slides_data, source_path)
        if output_pptx is None:
            return {}

        try:
            processed_urls = SlideRenderService.export_slide_images(
                str(output_pptx),
                courseware_id,
                output_folder="processed_slides",
            )
            debug_log(
                hypothesisId="H12",
                runId="pre-diagnose",
                location="image_processing_service:process_all_slides",
                message="Rendered translated slides from editable PPT",
                data={"courseware_id": courseware_id, "processed_count": len(processed_urls)},
            )
            return processed_urls
        except Exception as exc:
            debug_log(
                hypothesisId="H13",
                runId="pre-diagnose",
                location="image_processing_service:process_all_slides",
                message="Failed to render translated slides from editable PPT",
                data={"courseware_id": courseware_id, "error": str(exc)[:500]},
            )
            return {}
