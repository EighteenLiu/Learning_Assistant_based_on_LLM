from __future__ import annotations

from collections import Counter
from dataclasses import dataclass
from pathlib import Path
import re

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from app.debug_logger import debug_log


@dataclass
class ParsedSlide:
    slide_no: int
    title: str
    source_text: str
    notes: str
    source_layout: dict


class PPTParserService:
    ROW_TOLERANCE = 0.02
    PDF_SHORT_PHRASE_MAX_CHARS = 32
    PDF_SHORT_PHRASE_MAX_LINES = 2
    PDF_SHORT_PHRASE_MAX_WORDS = 6
    PDF_SHORT_PHRASE_REPEAT_THRESHOLD = 3
    PPT_WATERMARK_MIN_ROTATION = 20.0
    PPT_WATERMARK_MAX_ROTATION = 70.0
    PPT_MIN_TRANSLATABLE_LATIN_CHARS = 2
    SUPPORTED_FORMATS = (".pptx", ".ppt", ".pdf")
    PPT_FORMATS = (".pptx", ".ppt")
    PDF_FORMATS = (".pdf",)
    EXCLUDED_PLACEHOLDERS = {
        getattr(PP_PLACEHOLDER, "DATE", None),
        getattr(PP_PLACEHOLDER, "FOOTER", None),
        getattr(PP_PLACEHOLDER, "SLIDE_NUMBER", None),
        getattr(PP_PLACEHOLDER, "HEADER", None),
    }

    @staticmethod
    def _normalize_text(value: str) -> str:
        if not value:
            return ""
        lines = [line.strip() for line in str(value).replace("\xa0", " ").splitlines()]
        return "\n".join([line for line in lines if line])

    @staticmethod
    def _flatten_text(value: str) -> str:
        return re.sub(r"\s+", " ", PPTParserService._normalize_text(value)).strip()

    @staticmethod
    def _text_script_counts(value: str) -> dict[str, int]:
        flattened = PPTParserService._flatten_text(value)
        return {
            "latin": sum(1 for char in flattened if char.isascii() and char.isalpha()),
            "cjk": sum(1 for char in flattened if "\u4e00" <= char <= "\u9fff"),
            "digits": sum(1 for char in flattened if char.isdigit()),
        }

    @staticmethod
    def _has_translatable_latin_text(value: str) -> bool:
        flattened = PPTParserService._flatten_text(value)
        return bool(re.search(r"[A-Za-z]{2,}", flattened))

    @staticmethod
    def _has_min_latin_word_count(value: str, min_words: int = 3) -> bool:
        flattened = PPTParserService._flatten_text(value)
        words = re.findall(r"[A-Za-z]{2,}(?:[-'][A-Za-z]{2,})*", flattened)
        return len(words) >= max(int(min_words or 0), 1)

    @staticmethod
    def _is_mostly_cjk_text(value: str) -> bool:
        counts = PPTParserService._text_script_counts(value)
        return counts["cjk"] >= 2 and counts["cjk"] >= counts["latin"]

    @staticmethod
    def _is_numeric_or_symbolic_text(value: str) -> bool:
        counts = PPTParserService._text_script_counts(value)
        return counts["latin"] == 0 and counts["cjk"] == 0

    @staticmethod
    def _normalized_rotation(rotation_deg: object) -> float:
        try:
            rotation = abs(float(rotation_deg or 0.0)) % 360.0
        except Exception:
            return 0.0
        return min(rotation, 360.0 - rotation)

    @staticmethod
    def _is_likely_rotated_watermark(container: dict) -> bool:
        flattened = PPTParserService._flatten_text(str(container.get("text", "") or ""))
        if not flattened or bool(container.get("is_title")):
            return False
        rotation = PPTParserService._normalized_rotation(container.get("rotation_deg"))
        if rotation < PPTParserService.PPT_WATERMARK_MIN_ROTATION:
            return False
        if rotation > PPTParserService.PPT_WATERMARK_MAX_ROTATION:
            return False
        return PPTParserService._is_repeated_pdf_short_phrase(flattened)

    @staticmethod
    def should_translate_ppt_text_container(
        container: dict,
        repeated_short_phrase_fingerprints: set[str] | None = None,
    ) -> bool:
        if str(container.get("kind", "") or "") == "image_ocr":
            return True

        text = str(container.get("text", "") or "").strip()
        if not text:
            return False

        if repeated_short_phrase_fingerprints:
            fingerprint = PPTParserService._pdf_text_fingerprint(text)
            if fingerprint and fingerprint in repeated_short_phrase_fingerprints:
                return False

        if PPTParserService._is_mostly_cjk_text(text):
            return False
        if PPTParserService._is_numeric_or_symbolic_text(text):
            return False
        if PPTParserService._is_likely_rotated_watermark(container):
            return False
        return PPTParserService._has_translatable_latin_text(text)

    @staticmethod
    def filter_ppt_translation_containers(
        containers: list[dict],
        repeated_short_phrase_fingerprints: set[str] | None = None,
    ) -> list[dict]:
        filtered: list[dict] = []
        for container in containers:
            if str(container.get("kind", "") or "") == "image_ocr":
                filtered.append(dict(container))
                continue
            if PPTParserService.should_translate_ppt_text_container(
                container,
                repeated_short_phrase_fingerprints=repeated_short_phrase_fingerprints,
            ):
                filtered.append(dict(container))
        return filtered

    @staticmethod
    def collect_ppt_repeated_short_phrase_fingerprints(slide_layouts: list[dict]) -> set[str]:
        repeated_candidates: Counter[str] = Counter()
        for layout in slide_layouts:
            containers = layout.get("text_containers", []) if isinstance(layout, dict) else []
            seen_on_slide: set[str] = set()
            for container in containers if isinstance(containers, list) else []:
                if str(container.get("kind", "") or "") == "image_ocr":
                    continue
                text = str(container.get("text", "") or "").strip()
                if not PPTParserService._is_repeated_pdf_short_phrase(text):
                    continue
                fingerprint = PPTParserService._pdf_text_fingerprint(text)
                if fingerprint:
                    seen_on_slide.add(fingerprint)
            for fingerprint in seen_on_slide:
                repeated_candidates[fingerprint] += 1
        return {
            key
            for key, count in repeated_candidates.items()
            if count >= PPTParserService.PDF_SHORT_PHRASE_REPEAT_THRESHOLD
        }

    @staticmethod
    def sanitize_ppt_source_layout(
        source_layout: dict,
        repeated_short_phrase_fingerprints: set[str] | None = None,
    ) -> tuple[str, str, dict]:
        slide_width = int(source_layout.get("page_width", 1) or 1)
        slide_height = int(source_layout.get("page_height", 1) or 1)
        containers = source_layout.get("text_containers", []) or []
        filtered_containers = PPTParserService.filter_ppt_translation_containers(
            containers,
            repeated_short_phrase_fingerprints=repeated_short_phrase_fingerprints,
        )
        return PPTParserService._finalize_layout(filtered_containers, slide_width, slide_height)

    @staticmethod
    def should_keep_image_ocr_result(source_text: str, translated_text: str) -> bool:
        normalized_source = PPTParserService._normalize_text(source_text)
        normalized_translated = PPTParserService._normalize_text(translated_text)
        if not normalized_source or not normalized_translated:
            return False
        if PPTParserService._is_mostly_cjk_text(normalized_source):
            return False
        if PPTParserService._has_min_latin_word_count(normalized_source, 3):
            return True
        if PPTParserService._is_numeric_or_symbolic_text(normalized_source):
            return False
        return PPTParserService._has_translatable_latin_text(normalized_source)

    @staticmethod
    def _normalize_paragraph_text(paragraph) -> str:
        text = PPTParserService._normalize_text(getattr(paragraph, "text", ""))
        if not text:
            return ""

        level = int(getattr(paragraph, "level", 0) or 0)
        indent = "  " * min(level, 4)
        bullet = "- " if level > 0 else ""
        return f"{indent}{bullet}{text}".rstrip()

    @staticmethod
    def _shape_placeholder_priority(shape) -> int:
        if not getattr(shape, "is_placeholder", False):
            return 50

        try:
            placeholder_type = shape.placeholder_format.type
        except Exception:
            return 40

        if placeholder_type == PP_PLACEHOLDER.TITLE:
            return 0
        if placeholder_type == PP_PLACEHOLDER.CENTER_TITLE:
            return 1
        if placeholder_type in {PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}:
            return 5
        return 20

    @staticmethod
    def _is_title_shape(shape) -> bool:
        if bool(getattr(shape, "is_title", False)):
            return True
        if not getattr(shape, "is_placeholder", False):
            return False
        try:
            return shape.placeholder_format.type in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
        except Exception:
            return False

    @staticmethod
    def _shape_intersects_slide(shape, slide_width: int, slide_height: int) -> bool:
        left = int(getattr(shape, "left", 0) or 0)
        top = int(getattr(shape, "top", 0) or 0)
        width = int(getattr(shape, "width", 0) or 0)
        height = int(getattr(shape, "height", 0) or 0)
        if width <= 0 or height <= 0:
            return False
        right = left + width
        bottom = top + height
        return right > 0 and bottom > 0 and left < slide_width and top < slide_height

    @staticmethod
    def _should_include_shape(shape, slide_width: int, slide_height: int) -> bool:
        if not PPTParserService._shape_intersects_slide(shape, slide_width, slide_height):
            return False
        if getattr(shape, "is_placeholder", False):
            try:
                if shape.placeholder_format.type in PPTParserService.EXCLUDED_PLACEHOLDERS:
                    return False
            except Exception:
                pass
        return True

    @staticmethod
    def _shape_path(path: list[int]) -> str:
        return ".".join(str(item) for item in path)

    @staticmethod
    def _container_key(kind: str, shape_path: list[int], row_index: int | None = None, col_index: int | None = None) -> str:
        base = f"{kind}:{PPTParserService._shape_path(shape_path)}"
        if row_index is not None and col_index is not None:
            return f"{base}:{row_index}:{col_index}"
        return base

    @staticmethod
    def _append_container(
        containers: list[dict],
        *,
        text: str,
        paragraphs: list[str],
        kind: str,
        shape,
        shape_path: list[int],
        left: int,
        top: int,
        width: int,
        height: int,
        slide_width: int,
        slide_height: int,
        priority: int,
        is_title: bool,
        row_index: int | None = None,
        col_index: int | None = None,
        allow_empty_text: bool = False,
        rotation_deg: float = 0.0,
    ) -> None:
        normalized_text = PPTParserService._normalize_text(text)
        if not normalized_text and not allow_empty_text:
            return

        containers.append(
            {
                "kind": kind,
                "container_key": PPTParserService._container_key(kind, shape_path, row_index, col_index),
                "shape_id": int(getattr(shape, "shape_id", 0) or 0),
                "shape_name": str(getattr(shape, "name", "") or ""),
                "shape_path": PPTParserService._shape_path(shape_path),
                "paragraphs": paragraphs,
                "text": normalized_text,
                "x": left / slide_width,
                "y": top / slide_height,
                "w": width / slide_width,
                "h": height / slide_height,
                "is_title": bool(is_title),
                "priority": priority,
                "row_index": row_index,
                "col_index": col_index,
                "rotation_deg": float(rotation_deg or 0.0),
            }
        )

    @staticmethod
    def _extract_picture_container(shape, shape_path: list[int], slide_width: int, slide_height: int) -> list[dict]:
        left = int(getattr(shape, "left", 0) or 0)
        top = int(getattr(shape, "top", 0) or 0)
        width = int(getattr(shape, "width", 0) or 0)
        height = int(getattr(shape, "height", 0) or 0)
        if width <= 0 or height <= 0:
            return []

        containers: list[dict] = []
        PPTParserService._append_container(
            containers,
            text="",
            paragraphs=[],
            kind="image_ocr",
            shape=shape,
            shape_path=shape_path,
            left=left,
            top=top,
            width=width,
            height=height,
            slide_width=slide_width,
            slide_height=slide_height,
            priority=80,
            is_title=False,
            allow_empty_text=True,
            rotation_deg=float(getattr(shape, "rotation", 0.0) or 0.0),
        )
        return containers

    @staticmethod
    def _container_to_blocks(container: dict) -> list[dict]:
        paragraphs = [paragraph for paragraph in container.get("paragraphs", []) if paragraph]
        if not paragraphs:
            return []

        height = float(container.get("h", 0) or 0)
        paragraph_height = height / max(len(paragraphs), 1)
        blocks: list[dict] = []
        for index, paragraph_text in enumerate(paragraphs):
            blocks.append(
                {
                    "text": paragraph_text,
                    "x": container.get("x", 0),
                    "y": float(container.get("y", 0) or 0) + index * paragraph_height,
                    "w": container.get("w", 0),
                    "h": paragraph_height,
                    "is_title": bool(container.get("is_title", False) and index == 0),
                    "priority": int(container.get("priority", 99)) + index,
                    "paragraph_index": index,
                    "container_key": container.get("container_key"),
                    "shape_id": container.get("shape_id", 0),
                    "font_size_pt": container.get("font_size_pt"),
                    "font_name": container.get("font_name", ""),
                }
            )
        return blocks

    @staticmethod
    def _extract_font_profile(text_frame) -> tuple[str, float]:
        font_name = ""
        paragraph_sizes: list[float] = []

        for paragraph in getattr(text_frame, "paragraphs", []):
            paragraph_font = getattr(paragraph, "font", None)
            if paragraph_font is not None:
                if not font_name and getattr(paragraph_font, "name", None):
                    font_name = str(paragraph_font.name)
                if getattr(paragraph_font, "size", None):
                    try:
                        paragraph_sizes.append(float(paragraph_font.size.pt))
                    except Exception:
                        pass

            for run in getattr(paragraph, "runs", []):
                run_font = getattr(run, "font", None)
                if run_font is None:
                    continue
                if not font_name and getattr(run_font, "name", None):
                    font_name = str(run_font.name)
                if getattr(run_font, "size", None):
                    try:
                        paragraph_sizes.append(float(run_font.size.pt))
                    except Exception:
                        pass

        preferred_size = paragraph_sizes[0] if paragraph_sizes else 20.0
        if paragraph_sizes:
            preferred_size = sorted(paragraph_sizes)[len(paragraph_sizes) // 2]
        return font_name, preferred_size

    @staticmethod
    def _extract_text_frame_container(shape, shape_path: list[int], slide_width: int, slide_height: int) -> list[dict]:
        paragraphs = [PPTParserService._normalize_paragraph_text(paragraph) for paragraph in shape.text_frame.paragraphs]
        paragraphs = [paragraph for paragraph in paragraphs if paragraph]
        if not paragraphs:
            return []
        font_name, font_size_pt = PPTParserService._extract_font_profile(shape.text_frame)

        containers: list[dict] = []
        PPTParserService._append_container(
            containers,
            text="\n".join(paragraphs),
            paragraphs=paragraphs,
            kind="text_frame",
            shape=shape,
            shape_path=shape_path,
            left=int(getattr(shape, "left", 0) or 0),
            top=int(getattr(shape, "top", 0) or 0),
            width=int(getattr(shape, "width", 0) or 0),
            height=int(getattr(shape, "height", 0) or 0),
            slide_width=slide_width,
            slide_height=slide_height,
            priority=PPTParserService._shape_placeholder_priority(shape),
            is_title=PPTParserService._is_title_shape(shape),
            rotation_deg=float(getattr(shape, "rotation", 0.0) or 0.0),
        )
        if containers:
            containers[-1]["font_name"] = font_name
            containers[-1]["font_size_pt"] = font_size_pt
        return containers

    @staticmethod
    def _extract_table_containers(shape, shape_path: list[int], slide_width: int, slide_height: int) -> list[dict]:
        containers: list[dict] = []
        rows = len(shape.table.rows) or 1
        cols = len(shape.table.columns) or 1
        cell_width = int((getattr(shape, "width", 0) or 0) / cols)
        cell_height = int((getattr(shape, "height", 0) or 0) / rows)
        base_left = int(getattr(shape, "left", 0) or 0)
        base_top = int(getattr(shape, "top", 0) or 0)

        for row_index, row in enumerate(shape.table.rows):
            for col_index, cell in enumerate(row.cells):
                paragraphs = [
                    normalized
                    for normalized in (
                        PPTParserService._normalize_text(item) for item in str(getattr(cell, "text", "") or "").splitlines()
                    )
                    if normalized
                ]
                if not paragraphs:
                    continue
                PPTParserService._append_container(
                    containers,
                    text="\n".join(paragraphs),
                    paragraphs=paragraphs,
                    kind="table_cell",
                    shape=shape,
                    shape_path=shape_path,
                    left=base_left + col_index * cell_width,
                    top=base_top + row_index * cell_height,
                    width=cell_width,
                    height=cell_height,
                    slide_width=slide_width,
                    slide_height=slide_height,
                    priority=30 + row_index * cols + col_index,
                    is_title=False,
                    row_index=row_index,
                    col_index=col_index,
                    rotation_deg=float(getattr(shape, "rotation", 0.0) or 0.0),
                )
                if containers:
                    font_name, font_size_pt = PPTParserService._extract_font_profile(cell.text_frame)
                    containers[-1]["font_name"] = font_name
                    containers[-1]["font_size_pt"] = font_size_pt
        return containers

    @staticmethod
    def _extract_slide_containers(shapes, slide_width: int, slide_height: int, path_prefix: list[int] | None = None) -> list[dict]:
        path_prefix = path_prefix or []
        containers: list[dict] = []

        for index, shape in enumerate(shapes):
            shape_path = [*path_prefix, index]
            shape_type = getattr(shape, "shape_type", None)

            if not PPTParserService._should_include_shape(shape, slide_width, slide_height):
                continue

            if shape_type == MSO_SHAPE_TYPE.GROUP:
                containers.extend(
                    PPTParserService._extract_slide_containers(shape.shapes, slide_width, slide_height, shape_path)
                )
                continue

            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                containers.extend(
                    PPTParserService._extract_picture_container(shape, shape_path, slide_width, slide_height)
                )
                continue

            if getattr(shape, "has_table", False):
                containers.extend(
                    PPTParserService._extract_table_containers(shape, shape_path, slide_width, slide_height)
                )
                continue

            if getattr(shape, "has_text_frame", False):
                containers.extend(
                    PPTParserService._extract_text_frame_container(shape, shape_path, slide_width, slide_height)
                )

        return containers

    @staticmethod
    def _sort_blocks(blocks: list[dict]) -> list[dict]:
        if not blocks:
            return []
        return sorted(
            blocks,
            key=lambda block: (
                round(float(block.get("y", 0) or 0) / max(PPTParserService.ROW_TOLERANCE, 0.0001)),
                int(block.get("priority", 99)),
                round(float(block.get("x", 0) or 0), 4),
                round(float(block.get("h", 0) or 0), 4),
            ),
        )

    @staticmethod
    def _sort_containers(containers: list[dict]) -> list[dict]:
        if not containers:
            return []
        return sorted(
            containers,
            key=lambda item: (
                round(float(item.get("y", 0) or 0) / max(PPTParserService.ROW_TOLERANCE, 0.0001)),
                int(item.get("priority", 99)),
                round(float(item.get("x", 0) or 0), 4),
            ),
        )

    @staticmethod
    def _finalize_layout(containers: list[dict], slide_width: int, slide_height: int) -> tuple[str, str, dict]:
        sorted_containers = PPTParserService._sort_containers(containers)
        blocks: list[dict] = []
        title = ""
        source_texts: list[str] = []

        for container_index, container in enumerate(sorted_containers, start=1):
            container["container_id"] = container_index
            if container.get("is_title") and not title:
                title = str(container.get("text", "")).strip()
            if str(container.get("text", "")).strip():
                source_texts.append(str(container.get("text", "")).strip())
            blocks.extend(PPTParserService._container_to_blocks(container))

        sorted_blocks = PPTParserService._sort_blocks(blocks)
        for block_index, block in enumerate(sorted_blocks, start=1):
            block["block_id"] = block_index
            container_id = 0
            for container in sorted_containers:
                if container.get("container_key") == block.get("container_key"):
                    container_id = int(container.get("container_id") or 0)
                    break
            block["container_id"] = container_id
            block.pop("priority", None)
            block.pop("container_key", None)

        for container in sorted_containers:
            container.pop("priority", None)

        return (
            title,
            "\n".join([text for text in source_texts if text]).strip(),
            {
                "page_width": slide_width,
                "page_height": slide_height,
                "blocks": sorted_blocks,
                "text_containers": sorted_containers,
            },
        )

    @staticmethod
    def _detect_file_ext(file_path: str) -> str:
        return str(Path(file_path).suffix or "").strip().lower()

    @staticmethod
    def is_pdf_file(file_path: str) -> bool:
        return PPTParserService._detect_file_ext(file_path) in set(PPTParserService.PDF_FORMATS)

    @staticmethod
    def is_ppt_file(file_path: str) -> bool:
        return PPTParserService._detect_file_ext(file_path) in set(PPTParserService.PPT_FORMATS)

    @staticmethod
    def parse_courseware(file_path: str) -> list[ParsedSlide]:
        suffix = PPTParserService._detect_file_ext(file_path)
        if suffix in PPTParserService.PPT_FORMATS:
            return PPTParserService.parse_pptx(file_path)
        if suffix in PPTParserService.PDF_FORMATS:
            return PPTParserService.parse_pdf(file_path)
        raise ValueError(f"Only {', '.join(PPTParserService.SUPPORTED_FORMATS)} files are supported.")

    @staticmethod
    def _safe_float(value: object, default: float = 0.0) -> float:
        try:
            return float(value)
        except Exception:
            return default

    @staticmethod
    def _extract_pdf_text_block(raw_block: dict) -> tuple[list[str], str, float]:
        lines = raw_block.get("lines", []) if isinstance(raw_block, dict) else []
        paragraphs: list[str] = []
        font_name = ""
        font_sizes: list[float] = []

        for line in lines if isinstance(lines, list) else []:
            spans = line.get("spans", []) if isinstance(line, dict) else []
            text_segments: list[str] = []
            for span in spans if isinstance(spans, list) else []:
                segment = PPTParserService._normalize_text(str(span.get("text", "") or ""))
                if not segment:
                    continue
                text_segments.append(segment)
                if not font_name:
                    font_name = str(span.get("font", "") or "")
                size = PPTParserService._safe_float(span.get("size"), 0.0)
                if size > 0:
                    font_sizes.append(size)
            line_text = PPTParserService._normalize_text(" ".join(text_segments))
            if line_text:
                paragraphs.append(line_text)

        preferred_size = 14.0
        if font_sizes:
            sorted_sizes = sorted(font_sizes)
            preferred_size = sorted_sizes[len(sorted_sizes) // 2]
        return paragraphs, font_name, preferred_size

    @staticmethod
    def _pdf_text_fingerprint(text: str) -> str:
        normalized = PPTParserService._normalize_text(text)
        if not normalized:
            return ""
        normalized = re.sub(r"\s+", " ", normalized).strip().casefold()
        return normalized

    @staticmethod
    def _is_repeated_pdf_short_phrase(text: str) -> bool:
        normalized = PPTParserService._normalize_text(text)
        if not normalized:
            return False

        lines = [line.strip() for line in normalized.splitlines() if line.strip()]
        if len(lines) > PPTParserService.PDF_SHORT_PHRASE_MAX_LINES:
            return False

        flattened = re.sub(r"\s+", " ", normalized).strip()
        if len(flattened) > PPTParserService.PDF_SHORT_PHRASE_MAX_CHARS:
            return False

        words = [word for word in flattened.split(" ") if word]
        if len(words) > PPTParserService.PDF_SHORT_PHRASE_MAX_WORDS:
            return False

        meaningful_chars = [char for char in flattened if char.isalnum() or ("\u4e00" <= char <= "\u9fff")]
        return len(meaningful_chars) >= 2

    @staticmethod
    def dedupe_pdf_repeated_short_phrases(containers: list[dict]) -> list[dict]:
        if not containers:
            return []

        repeated_candidates: Counter[str] = Counter()
        for container in containers:
            if str(container.get("kind", "")) == "image_ocr":
                continue
            text = str(container.get("text", "")).strip()
            if not PPTParserService._is_repeated_pdf_short_phrase(text):
                continue
            fingerprint = PPTParserService._pdf_text_fingerprint(text)
            if fingerprint:
                repeated_candidates[fingerprint] += 1

        duplicate_keys = {
            key
            for key, count in repeated_candidates.items()
            if count >= PPTParserService.PDF_SHORT_PHRASE_REPEAT_THRESHOLD
        }
        if not duplicate_keys:
            return list(containers)

        deduped: list[dict] = []
        seen_short_phrases: set[str] = set()
        for container in containers:
            if str(container.get("kind", "")) == "image_ocr":
                deduped.append(container)
                continue

            fingerprint = PPTParserService._pdf_text_fingerprint(str(container.get("text", "")).strip())
            if fingerprint in duplicate_keys:
                if fingerprint in seen_short_phrases:
                    continue
                seen_short_phrases.add(fingerprint)
            deduped.append(container)

        return deduped

    @staticmethod
    def _mark_pdf_titles(containers: list[dict]) -> None:
        text_containers = [item for item in containers if str(item.get("kind", "")) != "image_ocr" and str(item.get("text", "")).strip()]
        if not text_containers:
            return
        if any(bool(item.get("is_title")) for item in text_containers):
            return

        font_sizes = [PPTParserService._safe_float(item.get("font_size_pt"), 0.0) for item in text_containers]
        max_font = max(font_sizes) if font_sizes else 0.0
        title_candidate = None
        for item in sorted(
            text_containers,
            key=lambda value: (
                -PPTParserService._safe_float(value.get("font_size_pt"), 0.0),
                PPTParserService._safe_float(value.get("y"), 0.0),
                PPTParserService._safe_float(value.get("x"), 0.0),
            ),
        ):
            font_size = PPTParserService._safe_float(item.get("font_size_pt"), 0.0)
            top_ratio = PPTParserService._safe_float(item.get("y"), 0.0)
            if top_ratio > 0.35:
                continue
            if font_size >= max(16.0, max_font * 0.88):
                title_candidate = item
                break

        if title_candidate is None:
            top_most = sorted(
                text_containers,
                key=lambda value: (
                    PPTParserService._safe_float(value.get("y"), 0.0),
                    PPTParserService._safe_float(value.get("x"), 0.0),
                ),
            )[0]
            if PPTParserService._safe_float(top_most.get("y"), 0.0) <= 0.2:
                title_candidate = top_most

        if title_candidate is not None:
            title_candidate["is_title"] = True

    @staticmethod
    def parse_pdf(file_path: str) -> list[ParsedSlide]:
        path = Path(file_path)
        if path.suffix.lower() not in set(PPTParserService.PDF_FORMATS):
            raise ValueError(f"Only {', '.join(PPTParserService.PDF_FORMATS)} files are supported.")

        try:
            import fitz  # type: ignore
        except Exception as exc:
            raise RuntimeError("PDF parsing requires PyMuPDF. Please install it with `pip install PyMuPDF`.") from exc

        document = None
        try:
            document = fitz.open(str(path.resolve()))
            parsed: list[ParsedSlide] = []
            total_deduped_containers = 0

            for index, page in enumerate(document, start=1):
                page_rect = page.rect
                page_width = max(PPTParserService._safe_float(page_rect.width, 1.0), 1.0)
                page_height = max(PPTParserService._safe_float(page_rect.height, 1.0), 1.0)
                text_payload = page.get_text("dict") or {}
                raw_blocks = text_payload.get("blocks", []) if isinstance(text_payload, dict) else []
                containers: list[dict] = []

                for block_index, raw_block in enumerate(raw_blocks if isinstance(raw_blocks, list) else [], start=1):
                    if not isinstance(raw_block, dict):
                        continue
                    bbox = raw_block.get("bbox")
                    if not (isinstance(bbox, (list, tuple)) and len(bbox) == 4):
                        continue

                    left = max(PPTParserService._safe_float(bbox[0], 0.0), 0.0)
                    top = max(PPTParserService._safe_float(bbox[1], 0.0), 0.0)
                    right = min(PPTParserService._safe_float(bbox[2], page_width), page_width)
                    bottom = min(PPTParserService._safe_float(bbox[3], page_height), page_height)
                    width = max(right - left, 0.0)
                    height = max(bottom - top, 0.0)
                    if width <= 0.5 or height <= 0.5:
                        continue

                    block_type = int(raw_block.get("type", 0) or 0)
                    container_base = {
                        "shape_id": 0,
                        "shape_name": "pdf_block",
                        "shape_path": f"pdf.{index}.{block_index}",
                        "x": left / page_width,
                        "y": top / page_height,
                        "w": width / page_width,
                        "h": height / page_height,
                        "priority": 40 + block_index,
                    }
                    if block_type == 1:
                        containers.append(
                            {
                                **container_base,
                                "kind": "image_ocr",
                                "container_key": f"image_ocr:pdf:{index}:{block_index}",
                                "paragraphs": [],
                                "text": "",
                                "is_title": False,
                            }
                        )
                        continue

                    if block_type != 0:
                        continue

                    paragraphs, font_name, font_size_pt = PPTParserService._extract_pdf_text_block(raw_block)
                    text = PPTParserService._normalize_text("\n".join(paragraphs))
                    if not text:
                        continue

                    containers.append(
                        {
                            **container_base,
                            "kind": "text_frame",
                            "container_key": f"text_frame:pdf:{index}:{block_index}",
                            "paragraphs": paragraphs,
                            "text": text,
                            "is_title": False,
                            "font_name": font_name,
                            "font_size_pt": font_size_pt,
                        }
                    )

                original_count = len(containers)
                containers = PPTParserService.dedupe_pdf_repeated_short_phrases(containers)
                total_deduped_containers += max(original_count - len(containers), 0)
                PPTParserService._mark_pdf_titles(containers)
                title, source_text, source_layout = PPTParserService._finalize_layout(
                    containers,
                    int(round(page_width)),
                    int(round(page_height)),
                )
                parsed.append(
                    ParsedSlide(
                        slide_no=index,
                        title=title,
                        source_text=source_text,
                        notes="",
                        source_layout=source_layout,
                    )
                )

            debug_log(
                hypothesisId="H4",
                runId="pre-diagnose",
                location="ppt_parser_service:parse_pdf",
                message="Parsed PDF pages",
                data={
                    "slide_count": len(parsed),
                    "deduped_container_count": total_deduped_containers,
                    "suffix": path.suffix.lower(),
                    "exists": path.exists(),
                    "is_file": path.is_file(),
                },
            )
            return parsed
        except Exception as exc:
            debug_log(
                hypothesisId="H4",
                runId="pre-diagnose",
                location="ppt_parser_service:parse_pdf",
                message="PDF parse failed",
                data={
                    "exc_type": type(exc).__name__,
                    "error": str(exc)[:500],
                    "suffix": path.suffix.lower(),
                    "exists": path.exists(),
                    "is_file": path.is_file(),
                },
            )
            raise
        finally:
            if document is not None:
                try:
                    document.close()
                except Exception:
                    pass

    @staticmethod
    def _convert_ppt_to_pptx(ppt_path: str) -> str:
        try:
            import pythoncom
            from win32com.client import DispatchEx

            pythoncom.CoInitialize()
            powerpoint = None
            presentation = None
            try:
                input_path = Path(ppt_path).resolve()
                output_path = input_path.with_suffix(".pptx")

                powerpoint = DispatchEx("PowerPoint.Application")
                try:
                    powerpoint.Visible = 0
                    powerpoint.DisplayAlerts = 0
                except Exception:
                    pass
                presentation = powerpoint.Presentations.Open(str(input_path), WithWindow=False)
                presentation.SaveAs(str(output_path), 24)
                return str(output_path)
            finally:
                if presentation is not None:
                    try:
                        presentation.Close()
                    except Exception:
                        pass
                if powerpoint is not None:
                    try:
                        powerpoint.Quit()
                    except Exception:
                        pass
                try:
                    pythoncom.CoUninitialize()
                except Exception:
                    pass
        except Exception as exc:
            debug_log(
                hypothesisId="H14",
                runId="pre-diagnose",
                location="ppt_parser_service:_convert_ppt_to_pptx",
                message="Failed to convert .ppt to .pptx",
                data={"error": str(exc)[:500]},
            )
            raise

    @staticmethod
    def ensure_editable_pptx(file_path: str) -> tuple[str, Path | None]:
        path = Path(file_path)
        if path.suffix.lower() == ".pptx":
            return str(path), None
        converted_path = Path(PPTParserService._convert_ppt_to_pptx(file_path))
        return str(converted_path), converted_path

    @staticmethod
    def parse_pptx(file_path: str) -> list[ParsedSlide]:
        path = Path(file_path)
        temp_pptx: Path | None = None
        try:
            supported_formats = list(PPTParserService.PPT_FORMATS)
            if path.suffix.lower() not in supported_formats:
                raise ValueError(f"Only {', '.join(supported_formats)} files are supported.")

            parse_path, temp_pptx = PPTParserService.ensure_editable_pptx(file_path)
            prs = Presentation(parse_path)
            parsed: list[ParsedSlide] = []
            slide_width = int(prs.slide_width or 1)
            slide_height = int(prs.slide_height or 1)

            for index, slide in enumerate(prs.slides, start=1):
                containers = PPTParserService._extract_slide_containers(slide.shapes, slide_width, slide_height)
                title, source_text, source_layout = PPTParserService._finalize_layout(
                    containers,
                    slide_width,
                    slide_height,
                )

                notes = ""
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = PPTParserService._normalize_text(slide.notes_slide.notes_text_frame.text or "")

                parsed.append(
                    ParsedSlide(
                        slide_no=index,
                        title=title,
                        source_text=source_text,
                        notes=notes,
                        source_layout=source_layout,
                    )
                )

            repeated_short_phrase_fingerprints = PPTParserService.collect_ppt_repeated_short_phrase_fingerprints(
                [slide.source_layout for slide in parsed]
            )
            for slide in parsed:
                title, source_text, source_layout = PPTParserService.sanitize_ppt_source_layout(
                    slide.source_layout,
                    repeated_short_phrase_fingerprints=repeated_short_phrase_fingerprints,
                )
                slide.title = title
                slide.source_text = source_text
                slide.source_layout = source_layout

            if temp_pptx:
                try:
                    temp_pptx.unlink(missing_ok=True)
                except Exception as cleanup_exc:
                    debug_log(
                        hypothesisId="H15",
                        runId="pre-diagnose",
                        location="ppt_parser_service:parse_pptx",
                        message="Failed to clean up temporary .pptx file",
                        data={"error": str(cleanup_exc)[:200]},
                    )

            debug_log(
                hypothesisId="H4",
                runId="pre-diagnose",
                location="ppt_parser_service:parse_pptx",
                message="Parsed PPTX slides",
                data={
                    "slide_count": len(parsed),
                    "suffix": path.suffix.lower(),
                    "exists": path.exists(),
                    "is_file": path.is_file(),
                },
            )
            return parsed
        except Exception as exc:
            if temp_pptx:
                try:
                    temp_pptx.unlink(missing_ok=True)
                except Exception as cleanup_exc:
                    debug_log(
                        hypothesisId="H15",
                        runId="pre-diagnose",
                        location="ppt_parser_service:parse_pptx",
                        message="Failed to clean up temporary .pptx file",
                        data={"error": str(cleanup_exc)[:200]},
                    )

            debug_log(
                hypothesisId="H4",
                runId="pre-diagnose",
                location="ppt_parser_service:parse_pptx",
                message="PPT parse failed",
                data={
                    "exc_type": type(exc).__name__,
                    "error": str(exc)[:500],
                    "suffix": path.suffix.lower(),
                    "exists": path.exists(),
                    "is_file": path.is_file(),
                },
            )
            raise

    @staticmethod
    def extract_courseware_title(file_path: str, parsed_slides: list[ParsedSlide] | None = None) -> str:
        path = Path(file_path)
        fallback_title = path.stem
        temp_pptx: Path | None = None

        try:
            if parsed_slides:
                first_slide_title = str(parsed_slides[0].title or "").strip()
                if first_slide_title:
                    return first_slide_title

            parse_path, temp_pptx = PPTParserService.ensure_editable_pptx(file_path)
            presentation = Presentation(parse_path)

            for slide in presentation.slides:
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False) and PPTParserService._is_title_shape(shape):
                        title_text = PPTParserService._normalize_text(shape.text_frame.text or "")
                        if title_text:
                            return title_text

            metadata_title = str(getattr(presentation.core_properties, "title", "") or "").strip()
            if metadata_title:
                return metadata_title

            for parsed_slide in parsed_slides or []:
                title_text = str(parsed_slide.title or "").strip()
                if title_text:
                    return title_text
        except Exception:
            pass
        finally:
            if temp_pptx:
                try:
                    temp_pptx.unlink(missing_ok=True)
                except Exception:
                    pass

        return fallback_title
