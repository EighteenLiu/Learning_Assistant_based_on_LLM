import json
import tempfile
from pathlib import Path
from unittest.mock import patch

from django.contrib.auth.models import User
from django.core.files.uploadedfile import SimpleUploadedFile
from django.test import override_settings
from pptx import Presentation
from rest_framework.test import APITestCase

from learning.models import Courseware, SlideContent
from learning.services.llm_client import LLMClientError


class InlineThread:
    def __init__(self, target=None, args=(), kwargs=None, daemon=None):
        self.target = target
        self.args = args
        self.kwargs = kwargs or {}

    def start(self):
        if self.target:
            self.target(*self.args, **self.kwargs)


@override_settings(MEDIA_ROOT=tempfile.gettempdir())
class AuthAndFlowApiTests(APITestCase):
    def setUp(self):
        self.user = User.objects.create_user(username="u1", password="pass123456")
        token_resp = self.client.post("/api/auth/login", {"username": "u1", "password": "pass123456"}, format="json")
        self.token = token_resp.data["access"]
        self.client.credentials(HTTP_AUTHORIZATION=f"Bearer {self.token}")

    def _build_pptx_upload(self):
        prs = Presentation()
        prs.core_properties.title = "Deck Metadata Title"
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Visible Deck Title"
        slide.placeholders[1].text = "English content for test"
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = Path(tmp.name)
        prs.save(str(tmp_path))
        with open(tmp_path, "rb") as file_obj:
            data = file_obj.read()
        tmp_path.unlink(missing_ok=True)
        return SimpleUploadedFile(
            "test.pptx",
            data,
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    def _build_named_pptx_upload(self, filename: str):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Visible Deck Title"
        slide.placeholders[1].text = "English content for test"
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = Path(tmp.name)
        prs.save(str(tmp_path))
        with open(tmp_path, "rb") as file_obj:
            data = file_obj.read()
        tmp_path.unlink(missing_ok=True)
        return SimpleUploadedFile(
            filename,
            data,
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    def test_auth_required(self):
        self.client.credentials()
        resp = self.client.get("/api/coursewares")
        self.assertEqual(resp.status_code, 401)

    def test_upload_validation(self):
        bad_file = SimpleUploadedFile("bad.txt", b"dummy", content_type="text/plain")
        resp = self.client.post("/api/coursewares/upload", {"file": bad_file}, format="multipart")
        self.assertEqual(resp.status_code, 400)

    @patch("learning.views.threading.Thread", InlineThread)
    @patch("learning.services.translation_service.ImageProcessingService.process_all_slides", return_value={})
    @patch("learning.services.vector_index_service.chromadb.PersistentClient")
    @patch("learning.services.llm_client.OpenAICompatibleClient.chat")
    def test_full_chain_success(self, mock_chat, _mock_chroma, _mock_process_images):
        def mock_chat_response(messages, **kwargs):
            last_message = messages[-1].content
            if "Source blocks" in last_message:
                return json.dumps([{"block_id": 1, "text": "翻译结果"}], ensure_ascii=False)
            if "chapter_summary" in last_message:
                return json.dumps(
                    {
                        "chapter_summary": "总结",
                        "key_points": ["重点"],
                        "term_pairs": [{"en": "LLM", "zh": "大语言模型"}],
                        "mind_map": {"title": "Root", "children": [{"title": "Branch", "children": []}]},
                    },
                    ensure_ascii=False,
                )
            return "翻译结果"

        mock_chat.side_effect = mock_chat_response

        file_obj = self._build_pptx_upload()
        upload_resp = self.client.post("/api/coursewares/upload", {"file": file_obj}, format="multipart")
        self.assertEqual(upload_resp.status_code, 201)
        cw_id = upload_resp.data["courseware_id"]
        self.assertTrue(str(upload_resp.data["title"]).startswith("test"))

        translate_resp = self.client.post(f"/api/coursewares/{cw_id}/translate")
        self.assertEqual(translate_resp.status_code, 202)

        status_resp = self.client.get(f"/api/coursewares/{cw_id}/status")
        self.assertEqual(status_resp.status_code, 200)
        self.assertEqual(status_resp.data["status"], "translated")
        self.assertEqual(status_resp.data["total_slides"], 1)
        self.assertEqual(status_resp.data["translated_slides"], 1)
        self.assertIn("translation_duration_seconds", status_resp.data)
        self.assertIsNotNone(status_resp.data["translation_duration_seconds"])
        self.assertGreaterEqual(int(status_resp.data["translation_duration_seconds"]), 0)

        slides_resp = self.client.get(f"/api/coursewares/{cw_id}/slides")
        self.assertEqual(slides_resp.status_code, 200)
        self.assertGreaterEqual(len(slides_resp.data), 1)
        self.assertIn("翻译结果", slides_resp.data[0]["translated_text"])

        with patch("learning.services.qa_service.QAService.ask", return_value=("答案", [{"slide_no": 1, "snippet": "xx"}])):
            qa_resp = self.client.post(f"/api/coursewares/{cw_id}/qa", {"question": "什么是 LLM?"}, format="json")
        self.assertEqual(qa_resp.status_code, 200)
        self.assertEqual(qa_resp.data["answer"], "答案")

        summary_resp = self.client.post(f"/api/coursewares/{cw_id}/summary", {}, format="json")
        self.assertEqual(summary_resp.status_code, 200)
        self.assertEqual(summary_resp.data["chapter_summary"], "总结")
        self.assertIn("mind_map", summary_resp.data)

        records_resp = self.client.get(f"/api/coursewares/{cw_id}/records")
        self.assertEqual(records_resp.status_code, 200)
        self.assertIn("qa_records", records_resp.data)

        list_resp = self.client.get("/api/coursewares")
        self.assertEqual(list_resp.status_code, 200)
        self.assertTrue(list_resp.data)
        self.assertIn("translation_duration_seconds", list_resp.data[0])

    def test_missing_courseware(self):
        resp = self.client.get("/api/coursewares/999/slides")
        self.assertEqual(resp.status_code, 404)

    def test_user_isolation(self):
        other = User.objects.create_user(username="u2", password="pass123456")
        courseware = Courseware.objects.create(owner=other, title="o", file="coursewares/o.pptx")
        SlideContent.objects.create(courseware=courseware, slide_no=1, source_text="x")
        resp = self.client.get(f"/api/coursewares/{courseware.id}/slides")
        self.assertEqual(resp.status_code, 404)

    @patch("learning.services.llm_client.OpenAICompatibleClient.chat", side_effect=LLMClientError("invalid_api_key"))
    def test_summary_fallback_when_llm_unavailable(self, _mock_chat):
        file_obj = self._build_pptx_upload()
        upload_resp = self.client.post("/api/coursewares/upload", {"file": file_obj}, format="multipart")
        self.assertEqual(upload_resp.status_code, 201)
        cw_id = upload_resp.data["courseware_id"]

        summary_resp = self.client.post(f"/api/coursewares/{cw_id}/summary", {}, format="json")
        self.assertEqual(summary_resp.status_code, 200)
        self.assertTrue(summary_resp.data["chapter_summary"])
        self.assertIn("本地总结草稿", summary_resp.data["chapter_summary"])
        self.assertIsInstance(summary_resp.data["key_points"], list)
        self.assertIn("mind_map", summary_resp.data)

    def test_upload_keeps_original_filename_and_uses_windows_style_suffix_for_duplicates(self):
        file1 = self._build_named_pptx_upload("测试(1).pptx")
        resp1 = self.client.post("/api/coursewares/upload", {"file": file1}, format="multipart")
        self.assertEqual(resp1.status_code, 201)
        self.assertEqual(resp1.data["title"], "测试(1)")

        file2 = self._build_named_pptx_upload("测试(1).pptx")
        resp2 = self.client.post("/api/coursewares/upload", {"file": file2}, format="multipart")
        self.assertEqual(resp2.status_code, 201)
        self.assertEqual(resp2.data["title"], "测试(1)(2)")

        file3 = self._build_named_pptx_upload("测试(1).pptx")
        resp3 = self.client.post("/api/coursewares/upload", {"file": file3}, format="multipart")
        self.assertEqual(resp3.status_code, 201)
        self.assertEqual(resp3.data["title"], "测试(1)(3)")
