import json
import tempfile
from pathlib import Path
from unittest.mock import patch

from django.contrib.auth.models import User
from django.test import TestCase, override_settings
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from learning.models import Courseware, SlideContent, TermDictionary, TranslationCache
from learning.services.image_processing_service import ImageProcessingService
from learning.services.ppt_parser_service import PPTParserService
from learning.services.slide_render_service import SlideRenderService
from learning.services.summary_service import SummaryService
from learning.services.translation_service import TranslationService
from learning.services.vector_index_service import VectorIndexService


class PPTParserServiceTests(TestCase):
    def test_parse_multi_slide_pptx(self):
        with tempfile.TemporaryDirectory() as td:
            file_path = Path(td) / "sample.pptx"
            prs = Presentation()
            slide1 = prs.slides.add_slide(prs.slide_layouts[1])
            slide1.shapes.title.text = "Intro"
            slide1.placeholders[1].text = "Large Language Models"
            slide2 = prs.slides.add_slide(prs.slide_layouts[1])
            slide2.shapes.title.text = "Second"
            slide2.placeholders[1].text = "Bilingual Learning"
            prs.save(file_path)

            slides = PPTParserService.parse_pptx(str(file_path))

        self.assertEqual(len(slides), 2)
        self.assertEqual(slides[0].slide_no, 1)
        self.assertIn("Large Language Models", slides[0].source_text)

    def test_parse_records_text_container_locators(self):
        with tempfile.TemporaryDirectory() as td:
            file_path = Path(td) / "shape-locators.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Intro"
            slide.placeholders[1].text = "First paragraph\nSecond paragraph"
            slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(22)
            prs.save(file_path)

            parsed = PPTParserService.parse_pptx(str(file_path))

        layout = parsed[0].source_layout
        self.assertIn("text_containers", layout)
        self.assertGreaterEqual(len(layout["text_containers"]), 1)
        first_container = layout["text_containers"][0]
        self.assertIn("shape_id", first_container)
        self.assertIn("shape_path", first_container)
        self.assertIn("paragraphs", first_container)
        self.assertIn("font_size_pt", first_container)
        self.assertGreater(first_container["font_size_pt"], 0)

    def test_parse_invalid_extension(self):
        with self.assertRaises(ValueError):
            PPTParserService.parse_pptx("invalid.txt")

    def test_extract_courseware_title_prefers_cover_slide_title(self):
        with tempfile.TemporaryDirectory() as td:
            file_path = Path(td) / "title-source.pptx"
            prs = Presentation()
            prs.core_properties.title = "Metadata Deck Title"
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Visible Slide Title"
            prs.save(file_path)

            parsed = PPTParserService.parse_pptx(str(file_path))
            title = PPTParserService.extract_courseware_title(str(file_path), parsed)

        self.assertEqual(title, "Visible Slide Title")

    def test_dedupe_pdf_repeated_short_phrases_keeps_only_one_short_repeat(self):
        containers = [
            {"kind": "text_frame", "text": "CONFIDENTIAL", "x": 0.1, "y": 0.1},
            {"kind": "text_frame", "text": "CONFIDENTIAL", "x": 0.3, "y": 0.2},
            {"kind": "text_frame", "text": "CONFIDENTIAL", "x": 0.5, "y": 0.3},
            {"kind": "text_frame", "text": "Main body content", "x": 0.1, "y": 0.5},
        ]

        deduped = PPTParserService.dedupe_pdf_repeated_short_phrases(containers)

        self.assertEqual(len(deduped), 2)
        self.assertEqual([item["text"] for item in deduped], ["CONFIDENTIAL", "Main body content"])


class PromptAndSummaryTests(TestCase):
    def setUp(self):
        self.user = User.objects.create_user(username="tester", password="pass123456")
        self.courseware = Courseware.objects.create(owner=self.user, title="c1", file="coursewares/fake.pptx")
        SlideContent.objects.create(
            courseware=self.courseware,
            slide_no=1,
            title="A",
            source_text="LLM helps translation\nAttention drives alignment",
            notes="",
            source_layout={
                "page_width": 1280,
                "page_height": 720,
                "blocks": [
                    {"block_id": 1, "text": "LLM helps translation", "x": 0.1, "y": 0.1, "w": 0.6, "h": 0.1},
                    {"block_id": 2, "text": "Attention drives alignment", "x": 0.1, "y": 0.3, "w": 0.7, "h": 0.1},
                ],
            },
        )

    def test_build_translated_layout_dedupes_existing_pdf_watermark_containers(self):
        self.courseware.file = "coursewares/fake.pdf"
        self.courseware.save(update_fields=["file"])
        slide = self.courseware.slides.first()
        slide.source_text = "CONFIDENTIAL\nCONFIDENTIAL\nCore lesson"
        slide.source_layout = {
            "page_width": 1280,
            "page_height": 720,
            "text_containers": [
                {
                    "container_id": 1,
                    "kind": "text_frame",
                    "paragraphs": ["CONFIDENTIAL"],
                    "text": "CONFIDENTIAL",
                    "x": 0.1,
                    "y": 0.1,
                    "w": 0.2,
                    "h": 0.05,
                },
                {
                    "container_id": 2,
                    "kind": "text_frame",
                    "paragraphs": ["CONFIDENTIAL"],
                    "text": "CONFIDENTIAL",
                    "x": 0.3,
                    "y": 0.2,
                    "w": 0.2,
                    "h": 0.05,
                },
                {
                    "container_id": 3,
                    "kind": "text_frame",
                    "paragraphs": ["CONFIDENTIAL"],
                    "text": "CONFIDENTIAL",
                    "x": 0.5,
                    "y": 0.3,
                    "w": 0.2,
                    "h": 0.05,
                },
                {
                    "container_id": 4,
                    "kind": "text_frame",
                    "paragraphs": ["Core lesson"],
                    "text": "Core lesson",
                    "x": 0.1,
                    "y": 0.5,
                    "w": 0.4,
                    "h": 0.08,
                },
            ],
            "blocks": [],
        }
        slide.save(update_fields=["source_text", "source_layout"])

        translator = TranslationService()
        with patch.object(
            TranslationService,
            "_translate_containers",
            side_effect=lambda containers, term_hint: {
                int(container["container_id"]): {
                    "text": f"ZH:{container['text']}",
                    "paragraphs": [f"ZH:{container['text']}"],
                }
                for container in containers
            },
        ), patch.object(TranslationService, "_translate_image_containers", return_value=({}, {})):
            translated_text, translated_layout, enhanced_source_text = translator._build_translated_layout(slide, "")

        self.assertEqual(
            [item["text"] for item in translated_layout["text_containers"]],
            ["ZH:CONFIDENTIAL", "ZH:Core lesson"],
        )
        self.assertEqual(enhanced_source_text, "CONFIDENTIAL\nCore lesson")
        self.assertEqual(translated_text, "ZH:CONFIDENTIAL\nZH:Core lesson")

    @patch("learning.services.translation_service.ImageProcessingService.process_all_slides", return_value={})
    @patch("learning.services.llm_client.OpenAICompatibleClient.chat")
    def test_translation_uses_batch_mapping_and_term_dictionary(self, mock_chat, _mock_process_images):
        mock_chat.return_value = json.dumps(
            [
                {"block_id": 1, "text": "大语言模型有助于翻译"},
                {"block_id": 2, "text": "注意力机制帮助内容对齐"},
            ],
            ensure_ascii=False,
        )
        TermDictionary.objects.create(source_term="LLM", target_term="大语言模型")

        TranslationService().translate_courseware(self.courseware)

        slide = self.courseware.slides.first()
        self.assertEqual(mock_chat.call_count, 1)
        self.assertEqual(slide.translated_layout["blocks"][0]["text"], "大语言模型有助于翻译")
        self.assertEqual(slide.translated_layout["blocks"][1]["text"], "注意力机制帮助内容对齐")
        self.assertEqual(
            slide.translated_text,
            "大语言模型有助于翻译\n注意力机制帮助内容对齐",
        )
        self.assertEqual(self.courseware.status, Courseware.STATUS_TRANSLATED)

        messages = mock_chat.call_args.args[0]
        self.assertIn("LLM => 大语言模型", messages[1].content)

    @patch("learning.services.llm_client.OpenAICompatibleClient.chat")
    def test_summary_json_parsing(self, mock_chat):
        mock_chat.return_value = json.dumps(
            {
                "chapter_summary": "章节概述",
                "key_points": ["重点一", "重点二"],
                "term_pairs": [{"en": "LLM", "zh": "大语言模型"}],
                "mind_map": {"title": "课程全景", "children": [{"title": "主题一", "children": []}]},
            },
            ensure_ascii=False,
        )

        chapter_summary, key_points, term_pairs, mind_map = SummaryService().generate(self.courseware)
        self.assertEqual(chapter_summary, "章节概述")
        self.assertEqual(key_points, ["重点一", "重点二"])
        self.assertEqual(term_pairs[0]["en"], "LLM")
        self.assertEqual(mind_map["title"], "课程全景")


    @patch("learning.services.llm_client.OpenAICompatibleClient.chat")
    def test_summary_json_parsing_with_markdown_fence(self, mock_chat):
        mock_chat.return_value = """```json
{
  "chapter_summary": "带代码块的摘要",
  "key_points": ["要点A", "要点B"],
  "term_pairs": [{"en": "RDA", "zh": "资源描述与检索"}],
  "mind_map": {"title": "主题根节点", "children": [{"title": "分支", "children": []}]}
}
```"""

        chapter_summary, key_points, term_pairs, mind_map = SummaryService().generate(self.courseware)
        self.assertEqual(chapter_summary, "带代码块的摘要")
        self.assertEqual(key_points, ["要点A", "要点B"])
        self.assertEqual(term_pairs[0]["en"], "RDA")
        self.assertEqual(mind_map["title"], "主题根节点")


class VectorIndexServiceTests(TestCase):
    @patch("learning.services.vector_index_service.chromadb.PersistentClient")
    def test_index_rebuild_and_query(self, mock_client_cls):
        mock_collection = mock_client_cls.return_value.get_or_create_collection.return_value
        mock_collection.query.return_value = {
            "documents": [["doc1"]],
            "metadatas": [[{"slide_no": 1}]],
            "distances": [[0.12]],
        }

        service = VectorIndexService()
        fake_courseware = type("CoursewareObj", (), {"id": 1, "owner_id": 1})()
        fake_slide = type("SlideObj", (), {"slide_no": 1, "translated_text": "doc1", "source_text": "doc1"})()

        service.rebuild_courseware_index(fake_courseware, [fake_slide])
        hits = service.query(1, "question")

        self.assertEqual(hits[0]["slide_no"], 1)


class SlideRenderServiceTests(TestCase):
    def test_image_sort_key_uses_numeric_suffix(self):
        paths = [Path("幻灯片10.PNG"), Path("幻灯片2.PNG"), Path("幻灯片1.PNG")]
        ordered = sorted(paths, key=SlideRenderService._image_sort_key)
        self.assertEqual([path.name for path in ordered], ["幻灯片1.PNG", "幻灯片2.PNG", "幻灯片10.PNG"])


class ImageProcessingServiceTests(TestCase):
    def test_set_text_frame_content_reduces_font_size_for_long_text(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2.2), Inches(1.2))
        text_frame = textbox.text_frame
        text_frame.text = "Short"
        text_frame.paragraphs[0].font.size = Pt(28)

        ImageProcessingService._set_text_frame_content(
            text_frame,
            ["这是一段明显更长的中文翻译文本，需要通过缩小字号才能完整显示在原始文本框范围内。"],
            width_emu=int(textbox.width),
            height_emu=int(textbox.height),
        )

        applied_size = None
        for paragraph in text_frame.paragraphs:
            if paragraph.font.size:
                applied_size = paragraph.font.size.pt
                break
            for run in paragraph.runs:
                if run.font.size:
                    applied_size = run.font.size.pt
                    break
        self.assertIsNotNone(applied_size)
        self.assertLess(applied_size, 28)

    def test_set_text_frame_content_prefers_original_font_hint_before_shrinking(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4.8), Inches(1.2))
        text_frame = textbox.text_frame
        text_frame.text = "Short"
        text_frame.paragraphs[0].font.size = Pt(26)

        ImageProcessingService._set_text_frame_content(
            text_frame,
            ["较短的中文翻译"],
            width_emu=int(textbox.width),
            height_emu=int(textbox.height),
            font_name_hint="Microsoft YaHei",
            font_size_hint=18,
        )

        applied_size = None
        for paragraph in text_frame.paragraphs:
            if paragraph.font.size:
                applied_size = paragraph.font.size.pt
                break
            for run in paragraph.runs:
                if run.font.size:
                    applied_size = run.font.size.pt
                    break
        self.assertIsNotNone(applied_size)
        self.assertLessEqual(applied_size, 18)

    def test_set_text_frame_content_forces_high_contrast_text_color(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2.6), Inches(1.2))
        text_frame = textbox.text_frame
        text_frame.text = "Invisible"
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        ImageProcessingService._set_text_frame_content(
            text_frame,
            ["翻译后的内容需要保持清晰可读。"],
            width_emu=int(textbox.width),
            height_emu=int(textbox.height),
            host=textbox,
            background_rgb=(255, 255, 255),
        )

        paragraph_color = text_frame.paragraphs[0].font.color.rgb
        self.assertIsNotNone(paragraph_color)
        self.assertNotEqual(tuple(paragraph_color), (255, 255, 255))

    def test_set_text_frame_content_uses_aggressive_fitting_for_dense_hangul_text(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2.1), Inches(1.0))
        text_frame = textbox.text_frame
        text_frame.text = "Short"
        text_frame.paragraphs[0].font.size = Pt(24)

        ImageProcessingService._set_text_frame_content(
            text_frame,
            ["이 페이지는 한국어 번역 이후에도 모든 내용을 빠짐없이 표시해야 하므로 매우 촘촘한 배치가 필요합니다."],
            width_emu=int(textbox.width),
            height_emu=int(textbox.height),
            host=textbox,
        )

        self.assertLessEqual(int(text_frame.margin_left), int(Pt(1)))
        self.assertIn(text_frame.auto_size, {None, MSO_AUTO_SIZE.NONE, MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE})


class TranslationPerformanceStrategyTests(TestCase):
    @override_settings(TRANSLATION_CHUNK_MAX_CONTAINERS=2, TRANSLATION_CHUNK_MAX_CHARS=20000)
    @patch("learning.services.llm_client.OpenAICompatibleClient.chat")
    def test_container_translation_uses_adaptive_chunking(self, mock_chat):
        def side_effect(messages, **kwargs):
            prompt = messages[-1].content
            payload_text = prompt.split("Source blocks / containers:\n", 1)[1]
            items = json.loads(payload_text)
            return json.dumps(
                [
                    {"container_id": item["container_id"], "text": f"T{item['container_id']}", "paragraphs": [f"T{item['container_id']}"]}
                    for item in items
                ],
                ensure_ascii=False,
            )

        mock_chat.side_effect = side_effect
        translator = TranslationService()
        containers = [
            {"container_id": 1, "text": "alpha", "paragraphs": ["alpha"], "is_title": False},
            {"container_id": 2, "text": "beta", "paragraphs": ["beta"], "is_title": False},
            {"container_id": 3, "text": "gamma", "paragraphs": ["gamma"], "is_title": False},
            {"container_id": 4, "text": "delta", "paragraphs": ["delta"], "is_title": False},
            {"container_id": 5, "text": "epsilon", "paragraphs": ["epsilon"], "is_title": False},
        ]

        translated = translator._translate_containers(containers, term_hint="")

        self.assertEqual(mock_chat.call_count, 3)
        self.assertEqual(translated[1]["text"], "T1")
        self.assertEqual(translated[5]["text"], "T5")

    @patch("learning.services.llm_client.OpenAICompatibleClient.chat")
    def test_container_translation_retries_structured_then_fallbacks_only_missing(self, mock_chat):
        mock_chat.side_effect = [
            "INVALID_JSON",
            json.dumps([{"container_id": 1, "text": "translated-1", "paragraphs": ["translated-1"]}], ensure_ascii=False),
            "translated-2",
        ]

        translator = TranslationService()
        containers = [
            {"container_id": 1, "text": "source-1", "paragraphs": ["source-1"], "is_title": False},
            {"container_id": 2, "text": "source-2", "paragraphs": ["source-2"], "is_title": False},
        ]

        translated = translator._translate_containers(containers, term_hint="")

        self.assertEqual(mock_chat.call_count, 3)
        self.assertEqual(translated[1]["text"], "translated-1")
        self.assertEqual(translated[2]["text"], "translated-2")

    @patch("learning.services.llm_client.OpenAICompatibleClient.chat")
    def test_text_translation_uses_cache(self, mock_chat):
        mock_chat.return_value = "cached-result"
        translator = TranslationService()

        first = translator._translate_text("repeat me", "", translation_type=TranslationService.CACHE_TYPE_SLIDE_TEXT)
        second = translator._translate_text("repeat me", "", translation_type=TranslationService.CACHE_TYPE_SLIDE_TEXT)

        self.assertEqual(first, "cached-result")
        self.assertEqual(second, "cached-result")
        self.assertEqual(mock_chat.call_count, 1)
        self.assertEqual(TranslationCache.objects.count(), 1)

    @patch("learning.services.llm_client.OpenAICompatibleClient.chat")
    def test_container_translation_hits_cache_on_second_run(self, mock_chat):
        mock_chat.return_value = json.dumps(
            [
                {"container_id": 1, "text": "cached-1", "paragraphs": ["cached-1"]},
                {"container_id": 2, "text": "cached-2", "paragraphs": ["cached-2"]},
            ],
            ensure_ascii=False,
        )

        containers = [
            {"container_id": 1, "text": "same-1", "paragraphs": ["same-1"], "is_title": False},
            {"container_id": 2, "text": "same-2", "paragraphs": ["same-2"], "is_title": False},
        ]

        first = TranslationService()._translate_containers(containers, term_hint="")
        self.assertEqual(first[1]["text"], "cached-1")
        self.assertEqual(mock_chat.call_count, 1)

        mock_chat.reset_mock()
        second = TranslationService()._translate_containers(containers, term_hint="")
        self.assertEqual(second[2]["text"], "cached-2")
        self.assertEqual(mock_chat.call_count, 0)

    def test_split_to_match_paragraphs_avoids_empty_padding_for_single_line_response(self):
        result = TranslationService._split_to_match_paragraphs(
            "统一译文",
            ["line1", "line2", "line3"],
            translated_paragraphs=[],
        )
        self.assertEqual(result, ["统一译文"])
